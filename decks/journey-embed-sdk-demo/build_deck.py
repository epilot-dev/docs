"""Build the Journey Embed SDK demo deck — 6-slide edition."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EPILOT_BLUE = RGBColor(0x00, 0x33, 0x99)
EPILOT_ACCENT = RGBColor(0x39, 0x7A, 0xFF)
EPILOT_DARK = RGBColor(0x14, 0x1B, 0x33)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFB)
WARN = RGBColor(0xC0, 0x39, 0x2B)
GOOD = RGBColor(0x22, 0x8B, 0x4F)
MUTED = RGBColor(0x60, 0x6A, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=EPILOT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=EPILOT_DARK,
                bullet="•", line_space=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        if isinstance(item, tuple):
            head, sub = item
            r1 = p.add_run()
            r1.text = f"{bullet}  {head}  "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = sub
            r2.font.size = Pt(size - 1)
            r2.font.color.rgb = MUTED
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = f"{bullet}  {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def add_code(slide, x, y, w, h, code, *, size=12):
    bg = add_rect(slide, x, y, w, h, RGBColor(0x1E, 0x24, 0x3A))
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xE6, 0xEC, 0xFF)
    return bg


def add_header_band(slide, kicker, title):
    add_rect(slide, 0, 0, SW, Inches(0.18), EPILOT_ACCENT)
    add_text(slide, Inches(0.55), Inches(0.32), Inches(12), Inches(0.4),
             kicker, size=12, bold=True, color=EPILOT_ACCENT)
    add_text(slide, Inches(0.55), Inches(0.62), Inches(12.3), Inches(0.8),
             title, size=30, bold=True, color=EPILOT_DARK)
    add_rect(slide, Inches(0.55), Inches(1.42), Inches(0.7), Emu(40000), EPILOT_ACCENT)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.55), Inches(7.10), Inches(6), Inches(0.3),
             "Journey Embed SDK · Demo · Sales · Marketing · CS",
             size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.10), Inches(1.4), Inches(0.3),
             f"{page_num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------------ Slide 1
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, EPILOT_DARK)
    add_rect(s, Inches(0.8), Inches(2.4), Inches(0.18), Inches(2.0), EPILOT_ACCENT)
    add_text(s, Inches(1.15), Inches(2.3), Inches(11), Inches(0.5),
             "DEMO · BETA", size=14, bold=True, color=EPILOT_ACCENT)
    add_text(s, Inches(1.15), Inches(2.7), Inches(12), Inches(1.4),
             "Journey Embed SDK", size=58, bold=True, color=WHITE)
    add_text(s, Inches(1.15), Inches(4.0), Inches(12), Inches(0.7),
             "Putting epilot Journeys on customer sites — the modern way",
             size=22, color=RGBColor(0xC8, 0xD3, 0xF5))
    add_text(s, Inches(1.15), Inches(5.5), Inches(12), Inches(0.5),
             "From the legacy __epilot script to a chainable SDK + Web Component",
             size=14, color=RGBColor(0x9D, 0xA9, 0xCC))
    add_text(s, Inches(1.15), Inches(6.8), Inches(12), Inches(0.4),
             "epilot · Journeys team", size=12, color=RGBColor(0x9D, 0xA9, 0xCC))


# ------------------------------------------------------------------ Slide 2
def slide_legacy():
    s = prs.slides.add_slide(BLANK)
    add_header_band(s, "TODAY", "Legacy __epilot embed — and where it fell short")
    add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.5),
             "One bundle.js, an iframe, and a global window.__epilot namespace.",
             size=14, color=MUTED)

    add_code(s, Inches(0.55), Inches(2.3), Inches(5.8), Inches(2.5),
        '<button onclick="__epilot.enterFullScreen(\'<id>\')">\n'
        '  Open Journey\n'
        '</button>\n\n'
        '<script\n'
        '  data-ep-mode="full-screen"\n'
        '  data-ep-journeyIds="<id>"\n'
        '  src=".../bundle.js"></script>',
        size=11)

    add_text(s, Inches(6.6), Inches(2.3), Inches(6.4), Inches(0.4),
             "Where it fell short", size=16, bold=True, color=WARN)
    pain = [
        ("Iframe-only", "no native HTML element, hard to integrate"),
        ("Global state", "window.__epilot — brittle in SPAs / React"),
        ("Init-time only", "options frozen after init"),
        ("Custom events", "__epilot.on() instead of standard listeners"),
        ("Manual sizing", "minHeight tuned by hand"),
        ("No TS / npm", "copy-paste snippets only"),
    ]
    add_bullets(s, Inches(6.6), Inches(2.75), Inches(6.4), Inches(4),
                pain, size=12.5, line_space=1.3, bullet="✕")

    add_rect(s, Inches(0.55), Inches(5.1), Inches(5.8), Inches(1.7), LIGHT_BG)
    add_text(s, Inches(0.75), Inches(5.2), Inches(5.5), Inches(0.4),
             "Net effect for the business", size=14, bold=True, color=EPILOT_DARK)
    add_bullets(s, Inches(0.75), Inches(5.6), Inches(5.5), Inches(1.2), [
        "Each integration becomes a one-off",
        "Sales/CS own the workarounds",
        "Slows down customer onboarding",
    ], size=11.5, line_space=1.2)


# ------------------------------------------------------------------ Slide 3
def slide_introducing():
    s = prs.slides.add_slide(BLANK)
    add_header_band(s, "INTRODUCING", "Journey Embed SDK (Beta)")
    add_text(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.5),
             "One chainable JavaScript API. Two rendering backends. Same options.",
             size=15, color=EPILOT_DARK)

    add_rect(s, Inches(0.55), Inches(2.4), Inches(5.9), Inches(2.7), LIGHT_BG)
    add_rect(s, Inches(0.55), Inches(2.4), Inches(5.9), Inches(0.5), EPILOT_BLUE)
    add_text(s, Inches(0.55), Inches(2.4), Inches(5.9), Inches(0.5),
             "  .asIframe()  —  Rewritten iframe engine",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.75), Inches(3.05), Inches(5.5), Inches(2.0), [
        "Replaces the legacy __epilot bundle",
        "Strong style isolation",
        "Multiple journeys on one page",
        "Auto-sizing — no minHeight tuning",
    ], size=12, line_space=1.25)

    add_rect(s, Inches(6.95), Inches(2.4), Inches(5.9), Inches(2.7), LIGHT_BG)
    add_rect(s, Inches(6.95), Inches(2.4), Inches(5.9), Inches(0.5), EPILOT_ACCENT)
    add_text(s, Inches(6.95), Inches(2.4), Inches(5.9), Inches(0.5),
             "  .asWebComponent()  —  <epilot-journey>",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.15), Inches(3.05), Inches(5.5), Inches(2.0), [
        "Native HTML element, Shadow DOM",
        "Better performance & accessibility",
        "Drops into CMS / no-code builders",
        "Recommended default",
    ], size=12, line_space=1.25)

    add_code(s, Inches(0.55), Inches(5.3), Inches(12.3), Inches(1.5),
        "$epilot.embed(id).asWebComponent().mode('inline')\n"
        "  .topBar(true).scrollToTop(true).contextData({ source: 'checkout' })\n"
        "  .dataInjectionOptions({ initialStepIndex: 1, ... }).append('#target')",
        size=12)


# ------------------------------------------------------------------ Slide 4
def slide_capabilities():
    s = prs.slides.add_slide(BLANK)
    add_header_band(s, "CAPABILITIES", "What the new SDK unlocks")

    cards = [
        ("Modes",
         ["Inline — in page flow",
          "Full-Screen — modal overlay",
          "Inline → Full-Screen on progress",
          "Launcher Journeys (auto)"]),
        ("Data injection",
         ["Prefill from host page data",
          "Start at any step",
          "Lock specific fields",
          "URL params + .contextData()"]),
        ("Standards events",
         ["window 'message' — no globals",
          "JOURNEY_LOADED, PAGE_VIEW",
          "FORM_EVENT, PROGRESS",
          "Plug into GA4 / GTM / Segment"]),
        ("Distribution",
         ["Sync CDN — drop one tag",
          "Async CDN with onReady()",
          "npm + full TypeScript types",
          "React refs + clean .remove()"]),
    ]
    for i, (title, points) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = Inches(0.55 + col * 6.4)
        y = Inches(1.85 + row * 2.55)
        add_rect(s, x, y, Inches(6.0), Inches(2.35), LIGHT_BG)
        add_rect(s, x, y, Inches(0.18), Inches(2.35), EPILOT_ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(5.5), Inches(0.45),
                 title, size=18, bold=True, color=EPILOT_DARK)
        add_bullets(s, x + Inches(0.3), y + Inches(0.65), Inches(5.5), Inches(1.6),
                    points, size=12.5, line_space=1.25)


# ------------------------------------------------------------------ Slide 5
def slide_configurator_video():
    s = prs.slides.add_slide(BLANK)
    add_header_band(s, "DEMO TIME", "New Embed Configurator — live walkthrough")
    add_rect(s, Inches(0.55), Inches(1.85), Inches(8.5), Inches(5.0), EPILOT_DARK)
    play = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                              Inches(4.4), Inches(3.6), Inches(0.9), Inches(1.5))
    play.rotation = 30
    play.fill.solid()
    play.fill.fore_color.rgb = EPILOT_ACCENT
    play.line.fill.background()
    add_text(s, Inches(0.55), Inches(6.0), Inches(8.5), Inches(0.7),
             "▶  Play video — Embed Configurator walkthrough",
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(9.3), Inches(1.85), Inches(3.7), Inches(0.5),
             "What you'll see", size=16, bold=True, color=EPILOT_BLUE)
    add_bullets(s, Inches(9.3), Inches(2.3), Inches(3.7), Inches(5), [
        "New options surfaced",
        "Backend toggle: WC / iframe",
        "Mode preview live",
        "Copy-paste snippet",
        "Data injection helpers",
        "Context data fields",
    ], size=12.5, line_space=1.3)
    add_text(s, Inches(9.3), Inches(6.0), Inches(3.7), Inches(0.7),
             "(Video to be inserted in demo)",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ------------------------------------------------------------------ Slide 6
def slide_value():
    s = prs.slides.add_slide(BLANK)
    add_header_band(s, "VALUE", "What this means for Sales · Marketing · CS")
    persona = [
        ("Sales", EPILOT_BLUE, [
            "“Drop on any site” — incl. CMS / no-code",
            "Modern stack credibility (npm + TS)",
            "Faster integrations → faster revenue",
            "Multi-journey pages enable bigger deals",
        ]),
        ("Marketing", EPILOT_ACCENT, [
            "Inline + Launcher reduce CTA friction",
            "Context data + URL params = clean attribution",
            "Standards events plug into GA4 / GTM",
            "Better page perf via Shadow DOM",
        ]),
        ("Customer Success", EPILOT_DARK, [
            "Configurator covers more cases out of the box",
            "Fewer custom-code requests",
            "Migration is non-breaking",
            "Self-serve docs + storybook",
        ]),
    ]
    for i, (name, color, items) in enumerate(persona):
        x = Inches(0.55 + i * 4.27)
        add_rect(s, x, Inches(1.85), Inches(4.05), Inches(0.6), color)
        add_text(s, x, Inches(1.85), Inches(4.05), Inches(0.6),
                 name, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, Inches(2.5), Inches(4.05), Inches(3.5), LIGHT_BG)
        add_bullets(s, x + Inches(0.2), Inches(2.65), Inches(3.7), Inches(3.3),
                    items, size=12, line_space=1.4)

    add_rect(s, Inches(0.55), Inches(6.2), Inches(12.3), Inches(0.6), EPILOT_DARK)
    add_text(s, Inches(0.75), Inches(6.2), Inches(12), Inches(0.6),
             "Docs: docs.epilot.io › Journeys › SDK   ·   Storybook: embed.journey.epilot.io/stories   ·   npm: @epilot/journey-embed-sdk",
             size=11.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


slide_title()
slide_legacy()
slide_introducing()
slide_capabilities()
slide_configurator_video()
slide_value()

total = len(prs.slides)
for i, sld in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    add_footer(sld, i, total)

out = "journey-embed-sdk-demo.pptx"
prs.save(out)
print(f"Saved {out} — {total} slides")
